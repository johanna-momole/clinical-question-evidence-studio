"""PowerPoint export for EvidenceBrief using python-pptx.

Generates a slide deck from a persisted EvidenceBrief without adding
new claims. Source identifiers appear in speaker notes.
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

from src.schemas.brief import EvidenceBrief, EvidenceSnapshot

_APP_VERSION = "1.0.0"

# Colour palette
_NAVY = RGBColor(0x1E, 0x3A, 0x5F)
_BLUE = RGBColor(0x2C, 0x52, 0x82)
_LIGHT_BLUE = RGBColor(0xEB, 0xF4, 0xFF)
_AMBER = RGBColor(0x92, 0x40, 0x0E)
_AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
_GRAY = RGBColor(0x71, 0x80, 0x96)
_RED = RGBColor(0xC0, 0x39, 0x2B)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_DARK = RGBColor(0x2D, 0x37, 0x48)


def _set_slide_background(slide: Any, rgb: RGBColor) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_title_text(tf: Any, text: str, size: int = 28, bold: bool = True, color: RGBColor = _NAVY) -> None:
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_body_text(tf: Any, lines: list[str], size: int = 14, color: RGBColor = _DARK) -> None:
    tf.clear()
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color


def _add_disclaimer_box(slide: Any, prs: Any, text: str) -> None:
    left, top = Cm(1.5), Cm(16.5)
    width, height = Cm(23), Cm(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"⚠ {text}"
    run.font.size = Pt(9)
    run.font.color.rgb = _AMBER
    p.alignment = PP_ALIGN.LEFT


def _new_slide(prs: Any, layout_idx: int = 6) -> Any:
    layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
    return prs.slides.add_slide(layout)  # type: ignore[attr-defined]


def _add_text_block(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = _DARK,
    wrap: bool = True,
) -> None:
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


SHORT_DISCLAIMER = (
    "Educational prototype — not clinically validated — does not provide medical advice."
)


def generate_pptx(
    brief: EvidenceBrief,
    snapshot: EvidenceSnapshot | None,
    qa_results: list[dict],
    review_history: list[dict],
    question_text: str = "",
    pico_summary: dict | None = None,
    phenotype_summary: dict | None = None,
    cohort_attrition: list[dict] | None = None,
) -> bytes:
    """Render a PPTX from a persisted EvidenceBrief. Returns raw bytes.

    No new claims are added during rendering.
    Source identifiers are placed in speaker notes.
    """
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)

    # ── Slide 1: Title ─────────────────────────────────────────────────────────
    slide = _new_slide(prs, 0)
    slide.shapes.title.text = "Evidence Brief"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = (
            f"Clinical Question-Evidence Studio\n{brief.brief_id} · v{brief.version}\n"
            f"Generation mode: {brief.generation_mode} · "
            f"Review: {brief.human_review_status}"
        )
    _add_disclaimer_box(slide, prs, SHORT_DISCLAIMER)
    slide.notes_slide.notes_text_frame.text = (
        f"Brief ID: {brief.brief_id}\n"
        f"Content hash: {brief.content_hash}\n"
        f"Snapshot hash: {brief.evidence_snapshot_hash}\n"
        f"Data origin: {brief.data_origin}"
    )

    # ── Slide 2: Disclaimer ────────────────────────────────────────────────────
    slide = _new_slide(prs, 6)
    _add_text_block(slide, "Required Disclaimer", 1, 0.8, 23, 1.5, size=22, bold=True, color=_NAVY)
    _add_text_block(slide, brief.disclaimer, 1, 2.5, 23, 10, size=13, color=_AMBER)
    _add_text_block(
        slide,
        f"Data-origin notice: {brief.data_notice or '—'}",
        1, 13, 23, 3, size=11, color=_GRAY,
    )
    slide.notes_slide.notes_text_frame.text = "This slide must remain in the deck."

    # ── Slide 3: Project and data-source overview ──────────────────────────────
    slide = _new_slide(prs, 6)
    _add_text_block(slide, "Project Overview", 1, 0.8, 23, 1.5, size=22, bold=True, color=_NAVY)
    overview_lines = [
        "Clinical Question-Evidence Studio — educational clinical informatics prototype",
        "",
        "Patient data:  Entirely synthetic (Synthea FHIR R4) — not real patient data",
        "External evidence:  " + (brief.data_origin.replace("_", " ").title()),
        "",
        "Evidence sources:  PubMed · ClinicalTrials.gov · CMS Coverage",
        "Pipeline:  Question → Phenotype → Cohort → Evidence → Brief → Export",
        "",
        f"App version:  {_APP_VERSION}",
    ]
    _add_text_block(slide, "\n".join(overview_lines), 1, 2.5, 23, 12, size=12, color=_DARK)
    _add_disclaimer_box(slide, prs, SHORT_DISCLAIMER)

    # ── Slide 4: PICO and phenotype ────────────────────────────────────────────
    slide = _new_slide(prs, 6)
    _add_text_block(slide, "PICO Framework & Phenotype", 1, 0.8, 23, 1.5, size=22, bold=True, color=_NAVY)
    if pico_summary:
        pico_text = (
            f"Population:    {pico_summary.get('population', '—')}\n"
            f"Intervention:  {pico_summary.get('intervention', '—')}\n"
            f"Comparator:    {pico_summary.get('comparator', '—') or 'None'}\n"
            f"Outcomes:      {', '.join(pico_summary.get('outcomes', [])) or '—'}"
        )
    else:
        pico_text = "PICO data not available."
    _add_text_block(slide, pico_text, 1, 2.5, 23, 6, size=13, color=_DARK)
    if phenotype_summary:
        pheno_text = (
            f"Phenotype:  {phenotype_summary.get('name', '—')} "
            f"v{phenotype_summary.get('version', '—')}\n"
            f"{phenotype_summary.get('description', '')}"
        )
    else:
        pheno_text = "Phenotype data not available."
    _add_text_block(slide, pheno_text, 1, 9, 23, 5, size=12, color=_GRAY)
    _add_disclaimer_box(slide, prs, SHORT_DISCLAIMER)

    # ── Slide 5: Synthetic cohort attrition ────────────────────────────────────
    slide = _new_slide(prs, 6)
    _add_text_block(
        slide,
        "Synthetic Cohort Attrition",
        1, 0.8, 23, 1.5, size=22, bold=True, color=_NAVY,
    )
    _add_text_block(
        slide,
        "⚠ Synthetic data only — does not represent real patient populations",
        1, 2.3, 23, 0.8, size=11, color=_RED,
    )
    if cohort_attrition:
        atr_lines = ["Step  |  Label  |  In  |  Excluded  |  Out"]
        atr_lines += [
            f"  {r.get('Step', '')}  |  {r.get('Label', '')}  |  "
            f"{r.get('Records In', '')}  |  {r.get('Excluded', '')}  |  "
            f"{r.get('Records Out', '')}"
            for r in cohort_attrition
        ]
        atr_text = "\n".join(atr_lines)
    else:
        atr_text = "No attrition data available."
    _add_text_block(slide, atr_text, 1, 3.3, 23, 11, size=11, color=_DARK)
    _add_disclaimer_box(slide, prs, SHORT_DISCLAIMER)
    slide.notes_slide.notes_text_frame.text = (
        "All numbers reflect synthetic Synthea data. "
        "No real patient data is used in this application."
    )

    # ── Slide 6: External evidence overview ────────────────────────────────────
    slide = _new_slide(prs, 6)
    _add_text_block(slide, "External Evidence Overview", 1, 0.8, 23, 1.5, size=22, bold=True, color=_NAVY)
    record_count = len(snapshot.records) if snapshot else 0
    sources: dict[str, int] = {}
    if snapshot:
        for rec in snapshot.records:
            sname = rec.source_name or "unknown"
            sources[sname] = sources.get(sname, 0) + 1
    source_lines = [f"{src}: {count} record(s)" for src, count in sorted(sources.items())]
    ev_text = (
        f"Total evidence records in snapshot: {record_count}\n\n"
        + ("\n".join(source_lines) or "No source breakdown available.")
        + f"\n\nSnapshot hash: {brief.evidence_snapshot_hash[:24]}…"
        + f"\nData origin: {brief.data_origin.replace('_', ' ').title()}"
    )
    _add_text_block(slide, ev_text, 1, 2.5, 23, 10, size=13, color=_DARK)
    _add_disclaimer_box(slide, prs, SHORT_DISCLAIMER)

    # ── Slides 7+: Key findings by dimension ──────────────────────────────────
    dims: dict[str, list[Any]] = {}
    for claim in brief.claims:
        if claim.claim_type in ("supported", "exploratory"):
            dims.setdefault(claim.dimension, []).append(claim)

    for dim, claims in dims.items():
        slide = _new_slide(prs, 6)
        _add_text_block(
            slide,
            dim.replace("_", " ").title() + " Findings",
            1, 0.8, 23, 1.5, size=20, bold=True, color=_BLUE,
        )
        lines = []
        note_lines = []
        for claim in claims[:6]:
            cit_nums = [c.citation_number for c in claim.citations]
            cite = " ".join(f"[{n}]" for n in cit_nums) if cit_nums else ""
            badge = f"[{claim.claim_type.upper()}]"
            lines.append(f"• {badge} {claim.text}{(' ' + cite) if cite else ''}")
            if claim.uncertainty_note:
                lines.append(f"  → {claim.uncertainty_note}")
            for sid in claim.source_ids:
                note_lines.append(f"Source: {sid}")
        _add_text_block(slide, "\n".join(lines), 1, 2.5, 23, 12, size=12, color=_DARK)
        _add_disclaimer_box(slide, prs, SHORT_DISCLAIMER)
        if note_lines:
            slide.notes_slide.notes_text_frame.text = "\n".join(note_lines)

    # ── Slide: Evidence gaps ───────────────────────────────────────────────────
    if brief.evidence_gaps:
        slide = _new_slide(prs, 6)
        _add_text_block(slide, "Evidence Gaps", 1, 0.8, 23, 1.5, size=22, bold=True, color=_NAVY)
        gap_lines = [f"• {g.description}" for g in brief.evidence_gaps]
        _add_text_block(slide, "\n".join(gap_lines), 1, 2.5, 23, 12, size=13, color=_DARK)
        _add_disclaimer_box(slide, prs, SHORT_DISCLAIMER)

    # ── Slide: Limitations and QA ─────────────────────────────────────────────
    slide = _new_slide(prs, 6)
    _add_text_block(slide, "Limitations & QA", 1, 0.8, 23, 1.5, size=22, bold=True, color=_NAVY)
    lim_lines = [f"• {lim}" for lim in (brief.limitations or [])][:8]
    _add_text_block(
        slide,
        "Limitations:\n" + ("\n".join(lim_lines) or "None recorded."),
        1, 2.5, 11.5, 11, size=11, color=_DARK,
    )
    critical = sum(
        1 for r in qa_results if r.get("severity") == "critical" and r.get("status") == "failed"
    )
    warnings = sum(1 for r in qa_results if r.get("status") == "warning")
    qa_summary = (
        f"QA checks: {len(qa_results)} total\n"
        f"Critical failures: {critical}\n"
        f"Warnings: {warnings}\n\n"
        f"Review status: {brief.human_review_status}"
    )
    _add_text_block(slide, "QA Summary:\n" + qa_summary, 13, 2.5, 11.5, 11, size=11, color=_DARK)
    _add_disclaimer_box(slide, prs, SHORT_DISCLAIMER)

    # ── Slide: Provenance ─────────────────────────────────────────────────────
    slide = _new_slide(prs, 6)
    _add_text_block(slide, "Provenance & Traceability", 1, 0.8, 23, 1.5, size=22, bold=True, color=_NAVY)
    prov_lines = [
        f"Brief ID:          {brief.brief_id}",
        f"Version:           {brief.version}",
        f"Content hash:      {brief.content_hash[:32]}…",
        f"Snapshot hash:     {brief.evidence_snapshot_hash[:32]}…",
        f"Generation mode:   {brief.generation_mode}",
        f"Data origin:       {brief.data_origin}",
        f"Review status:     {brief.human_review_status}",
        f"Evidence records:  {len(snapshot.records) if snapshot else '—'}",
    ]
    _add_text_block(slide, "\n".join(prov_lines), 1, 2.5, 23, 10, size=12, color=_DARK)
    _add_disclaimer_box(slide, prs, SHORT_DISCLAIMER)
    slide.notes_slide.notes_text_frame.text = (
        f"Evidence run ID: {brief.evidence_run_id}\n"
        f"Snapshot ID: {brief.evidence_snapshot_id}\n"
        f"App version: {_APP_VERSION}"
    )

    # ── Final slide: Disclaimer ────────────────────────────────────────────────
    slide = _new_slide(prs, 6)
    _add_text_block(slide, "Disclaimer", 1, 0.8, 23, 1.5, size=22, bold=True, color=_NAVY)
    _add_text_block(slide, brief.disclaimer, 1, 2.5, 23, 12, size=14, color=_AMBER)
    slide.notes_slide.notes_text_frame.text = "Required disclaimer — must appear in all distributions."

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
